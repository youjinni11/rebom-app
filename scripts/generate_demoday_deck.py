#!/usr/bin/env python3
"""리봄 데모데이 피피티 생성 스크립트"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 브랜드 컬러
FOREST = RGBColor(0x1A, 0x47, 0x3A)
GOLD = RGBColor(0xC9, 0xA9, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xF8)

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "rebom-demoday-deck.pptx"


def set_run(run, size=18, bold=False, color=DARK, font_name="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_title_bar(slide, title: str, subtitle: Optional[str] = None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(12), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        set_run(sr, size=13, color=GOLD)


def add_bullets(slide, items: list[str], left=0.65, top=1.45, width=12.0, height=5.5, size=17, spacing=1.15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(size * spacing * 0.35)
        p.level = 0
        r = p.add_run()
        r.text = item
        set_run(r, size=size, color=DARK)


def add_two_column(slide, left_title, left_items, right_title, right_items):
    # Left column
    lt = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(5.9), Inches(0.4))
    ltf = lt.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = left_title
    set_run(lr, size=20, bold=True, color=FOREST)

    add_bullets(slide, left_items, left=0.65, top=1.85, width=5.9, height=4.8, size=15)

    # Right column
    rt = slide.shapes.add_textbox(Inches(6.85), Inches(1.35), Inches(5.9), Inches(0.4))
    rtf = rt.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = right_title
    set_run(rr, size=20, bold=True, color=FOREST)

    add_bullets(slide, right_items, left=6.85, top=1.85, width=5.9, height=4.8, size=15)


def add_metric_row(slide, metrics: list[tuple[str, str]], top=1.35):
    col_w = 12.0 / len(metrics)
    for i, (label, value) in enumerate(metrics):
        x = 0.65 + i * col_w
        card = slide.shapes.add_shape(1, Inches(x), Inches(top), Inches(col_w - 0.15), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = GOLD

        vbox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 0.15), Inches(col_w - 0.45), Inches(0.55))
        vtf = vbox.text_frame
        vtf.clear()
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vr = vp.add_run()
        vr.text = value
        set_run(vr, size=22, bold=True, color=FOREST)

        lbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(top + 0.75), Inches(col_w - 0.35), Inches(0.5))
        ltf = lbox.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = label
        set_run(lr, size=12, color=GRAY)


def add_footer_note(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(6.85), Inches(12), Inches(0.45))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size=11, color=GRAY)


def slide_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = FOREST
    bg.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.2))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "리봄"
    set_run(r, size=54, bold=True, color=WHITE)

    sbox = slide.shapes.add_textbox(Inches(0.9), Inches(3.35), Inches(11.5), Inches(1.0))
    stf = sbox.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "검증된 5060만을 위한 프리미엄 비대면 결정사 멤버십"
    set_run(sr, size=24, color=GOLD)

    mbox = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.8))
    mtf = mbox.text_frame
    mtf.clear()
    mp = mtf.paragraphs[0]
    mr = mp.add_run()
    mr.text = "데모데이 발표 | 5분"
    set_run(mr, size=16, color=WHITE)


def add_speaker_notes(slide, text: str):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)


def slide_content(prs: Presentation, title: str, subtitle: Optional[str], bullets: list[str], footer: Optional[str] = None, notes: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title, subtitle)
    add_bullets(slide, bullets)
    if footer:
        add_footer_note(slide, footer)
    if notes:
        add_speaker_notes(slide, notes)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 타이틀
    slide_title(prs)

    # 2. 문제
    slide_content(
        prs,
        "문제 — 5060의 만남은 왜 막혀 있나",
        "첫인상: 불편·불신",
        [
            "50·60대는 새 만남을 원하지만, 믿고 쓸 수 있는 방법이 거의 없다",
            "데이팅 앱: 연령·UX·사기 불신 — 타깃에 맞지 않음",
            "소개·모임: 검증 없음, 품질 편차 큼, 부담스러운 자리",
            "결정사: 가입비 400~500만 원, 대면·결혼 중심 — 부담 과다",
            "혼자 사는 5060 1인 가구는 늘지만, 그에 맞는 '검증된 만남' 인프라는 비어 있다",
        ],
        notes="[25초] 5060은 만남을 원하지만 믿을 방법이 없다는 문제로 시작. 감성보다 불신·막힘 강조.",
    )

    # 3. 기존 대안
    slide_content(
        prs,
        "기존 대안의 한계",
        None,
        [
            "데이팅 앱 → 스와이프·맞관심 구조, 5060에게 낯설고 불안",
            "소개팅·모임 → 신원 검증 없음, 진지한 만남 니즈와 결이 다름",
            "결정사 → 고가·결혼 목적, 리봄 타깃(동반·검증된 만남)과 다름",
            "공통 공백: 미혼 인증 + 프리미엄 + 비대면 + 5060 전용 UX",
        ],
        footer="리봄은 '데이팅 앱'이 아니라 '검증된 멤버십' 포지션",
    )

    # 4. 리봄 정의
    slide_content(
        prs,
        "리봄 — 한 문장 정의",
        "중간발표 질문 ① 선제 답변",
        [
            "미혼 인증을 통과한 50·60대만 입장하는 프리미엄 비대면 결정사형 멤버십",
            "맞관심 없음 → 리봄이 월 2회 직접 맞선 배정 (초기 수동 → AI+팀 검수)",
            "결혼이 목적이 아님 — 검증된 동반·만남, 회원이 관계의 깊이를 선택",
            "5단계 검증(전화·얼굴·미혼·프로필·가치관) + 만남 당일까지 앱에서 일정 조율",
        ],
    )

    # 5. 수요 검증
    slide_content(
        prs,
        "수요 검증",
        "정성 인터뷰 60명 · 타깃 20명",
        [
            "심층 인터뷰 60명(실제 타깃 약 20명) — 설문이 아닌 직접 대화로 니즈 확인",
            "공통 니즈: '검증된 사람'과의 만남, 사기·허위 프로필에 대한 불안",
            "가격: 월 20만 원 멤버십 적정 수준 확인 (팀 인터뷰·사전 조사)",
            "피봇 3회(데이팅앱→연정사→비대면 결정사) 후 현재 모델에 수렴",
            "※ 외부 설문 '연애 필요 72.8%'는 시장 분위기 참고용 — SAM/SOM 산정에 미사용",
        ],
    )

    # 6. 5060 + 모바일
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "5060은 이미 모바일을 쓴다", "중간발표 질문 ③ 선제 답변")
    add_metric_row(
        slide,
        [
            ("유튜브 55세+ 비율", "77.9%"),
            ("모바일 시청", "94.2%"),
            ("인스타 비팔로워 도달", "90.6%"),
        ],
        top=1.55,
    )
    add_bullets(
        slide,
        [
            "유튜브(28일): 조회 6,700 · 영상 7개 — 타깃 연령이 실제로 콘텐츠를 소비",
            "인스타(30일): 조회 1,004 · 도달 806 — 알고리즘이 5060에게 노출",
            "5060 전용 UX(큰 글씨·단순 구조) + 이미 익숙한 스마트폰 = 앱 진입 장벽 낮음",
        ],
        top=3.2,
        size=16,
    )

    # 7. 제품 데모
    slide_content(
        prs,
        "제품 데모",
        "GIF · 스크린샷 삽입",
        [
            "[여기에 데모 GIF/스크린샷 삽입]",
            "온보딩 → 5단계 검증(미혼 인증) → 가치관·라이프스타일 큐레이션",
            "홈: 오늘의 추천 / 매칭·일정 / 내 정보",
            "MVP 구현 완료: 지역 매칭, Codef 미혼 검증, 프로필·추천 파이프라인",
        ],
        footer="데모 영상: docs/demo-*.md 참고 · rebom-libom.vercel.app",
    )

    # 8. TAM — 전용 페이지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "TAM — Total Addressable Market", "전 세계 시니어 매칭 시장")
    add_metric_row(
        slide,
        [
            ("2025 시장 규모", "~18억 달러"),
            ("2034 전망", "~41억 달러"),
            ("CAGR", "9.5%"),
        ],
        top=1.35,
    )
    add_two_column(
        slide,
        "추정 근거",
        [
            "출처: DataIntelo, Global Elderly Dating Platform Market Report (2025)",
            "정의: 50세+ 온라인 만남·매칭 플랫폼 시장 전체 매출",
            "리봄 포지션: 대량 데이팅이 아닌 프리미엄·검증 세그먼트",
            "성장 동력: 고령 인구 증가 + 시니어 온라인 만남 수요 확대",
        ],
        "실질적 도달 경로",
        [
            "Phase 1 (1~2년): 한국에서 모델 검증 — 강남 300 → 2,000명",
            "Phase 2 (2~3년): 수도권·전국 확대 + AI 맞선 본격화",
            "Phase 3 (4년+): 동아시아(일본·대만) → 영어권 시니어 프리미엄",
            "한국 = 검증 센터, 해외 = 검증된 모델 라이선스·제휴 확장",
        ],
    )
    add_footer_note(slide, "Phase 2+ 익명·집계 인사이트 B2B 제휴 검토 (TAM 수치 미포함)")
    add_speaker_notes(slide, "[35초] TAM은 글로벌 시니어 매칭 약 18억 달러. 리봄은 프리미엄·검증 세그먼트. 한국 검증 후 동아시아·영어권으로 확장.")

    # 9. SAM — 전용 페이지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "SAM — Serviceable Addressable Market", "한국에서 지금 손댈 수 있는 시장")
    add_metric_row(
        slide,
        [
            ("50~69세 1인가구", "256만"),
            ("스마트폰 보유", "248만"),
            ("SAM 최종", "≈248만 명"),
        ],
        top=1.35,
    )
    add_two_column(
        slide,
        "추정 근거 (출처 1:1)",
        [
            "STEP1 256만: 통계청 《2024 통계로 보는 1인가구》2023 — 782.9만×(50대15.4%+60대17.3%)",
            "STEP2 248만: 방통위 《2024 방송매체 이용행태조사》60대 스마트폰 보유율 96.9%",
            "STEP3 삭제: '연애 필요 72.8%'는 태도 설문 → 유료 의향으로 해석 불가",
            "만남 니즈·가격 적정성: 팀 인터뷰 60명(타깃 20명)으로 별도 검증",
        ],
        "실질적 도달 경로",
        [
            "1년차 집중: 수도권 1인가구 비중 높은 서울·경기 (전국 SAM의 핵심 밀도권)",
            "해변(headland): 강남·서초 300명 오픈 → 맞선·BM 검증",
            "확장: 서울 → 수도권 광역시 (부산·대구 등은 Phase 2)",
            "유입: 설명회·인터뷰 네트워크·SNS — 결혼 앱이 아닌 '검증 멤버십' 메시지",
        ],
    )
    add_footer_note(slide, "SAM = 구조적 풀(인구) · 유료 전환은 SOM·인터뷰에서 검증")
    add_speaker_notes(slide, "[35초] SAM 248만. 256만 1인가구×스마트폰. 72.8%는 태도 설문이라 SAM에 곱하지 않음. 강남 300이 첫 섬.")

    # 10. SOM — 전용 페이지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "SOM — Serviceable Obtainable Market", "1년 안에 실제로 가져올 몫")
    add_metric_row(
        slide,
        [
            ("1년 유료 회원", "2,000명"),
            ("SAM 대비", "0.08%"),
            ("누적 매출(보수)", "~20억 원"),
        ],
        top=1.35,
    )
    add_two_column(
        slide,
        "추정 근거",
        [
            "SAM 역산: 248만 × 0.081% ≈ 2,000명 (프리미엄은 낮은 침투율이 현실적)",
            "실행 역산: Q1 강남·서초 300명(M4) → Q4 누적 2,000명",
            "매출: 2,000명 × LTV 100만(입장 20만+월20만×4개월) ≈ 20억 누적",
            "1년 현금흐름 보수: 순차 가입 반영 시 약 14억 원",
        ],
        "실질적 도달 경로",
        [
            "Q0: 사전관심 100명 — 설명회 3회·인터뷰 소개·SNS",
            "Q1: 300명 — 강남·서초 설명회 6회, 남녀 각 150 균형",
            "Q2~Q3: 300→1,300 — 서울 확대, 추천 프로그램, AI 배정 전환(1,000명+)",
            "Q4: 2,000명 — 추천 비중 25%+, 유료 광고 축소",
        ],
    )
    add_footer_note(slide, "운영 KPI: 맞선 월 2회 배정 · 만남 일정 확정률 목표 60%")
    add_speaker_notes(slide, "[25초] SOM 2,000명은 SAM의 0.08%. M4에 강남 300, 1년 누적 매출 보수 20억. 설명회·추천으로 도달.")

    # 11. BM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "비즈니스 모델", "중간발표 질문 ② 선제 답변")
    add_metric_row(
        slide,
        [
            ("입장권 (1회)", "20만 원"),
            ("월 멤버십", "20만 원"),
            ("LTV (보수)", "100만 원"),
        ],
        top=1.35,
    )
    add_bullets(
        slide,
        [
            "월 20만 원 = 맞선 월 2회 직접 배정 + 검증 풀 유지 (맞관심 구조 없음)",
            "첫 만남 1회 무료 · 교제 중 홀딩 최대 3개월(월 멤버십 면제)",
            "왜 구독? → 결정사 대비 1/20 가격대, 검증·큐레이션·일정 인프라 지속 제공",
            "검증 1인 변동비 약 750~1,540원 (입장료의 ~1%) · 월 고정비 15~30만 원",
            "LTV/CAC: 실질 CAC 10~15만 목표, LTV 100만 → 약 6~10배",
        ],
        top=3.05,
        size=16,
    )

    # 12. KPI + 마케팅 + 로드맵
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "실행 계획 · KPI · 마케팅", "Year 1 기준안")
    add_two_column(
        slide,
        "12개월 KPI",
        [
            "유료·검증 회원: 2,000명 (M4 강남 300명)",
            "런칭 전 사전관심: 100명",
            "만남 일정 확정률: 60%",
            "남녀 비율: 각 45~55% 유지",
            "분기: Q1 300 → Q2 700 → Q3 1,300 → Q4 2,000",
        ],
        "마케팅 3,500만 원 (B안)",
        [
            "설명회·직접 모집 34% · 유료 광고 29% · 추천·오가닉 37%",
            "Q0 300만 → Q1 800만 → Q2~Q3 1,900만 → Q4 500만",
            "유료 광고 CAC 상한 22만 원 — 초과 채널 즉시 중단",
            "핵심 채널: 설명회 > 추천 > SNS > 제한적 유료 광고",
            "변호사 법률 리스크 검토 완료 · MVP 출시 준비",
        ],
    )
    add_footer_note(slide, "로드맵: Y1 검증(2,000) → Y2~3 전국 5,000+ AI맞선 → Y4+ 동아시아")

    # 13. 팀 + 마무리
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "팀 · 마무리", None)
    add_two_column(
        slide,
        "팀",
        [
            "송유진 — 기획, BM 설계, 법률조사, 수요조사",
            "김소정 — 자료 제작, 수요조사, 인터뷰",
            "배준수 — 시장조사, 인터뷰",
            "3명이 직접 60명 인터뷰하고 만든 서비스",
        ],
        "한 줄 요약",
        [
            "5060은 만남을 원하지만 믿을 방법이 없다",
            "리봄 = 검증된 5060만의 프리미엄 비대면 결정사 멤버십",
            "한국 248만 SAM 중 1년 2,000명으로 검증 → 아시아·글로벌 확장",
            "다시, 봄이 옵니다 — 리봄",
        ],
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
