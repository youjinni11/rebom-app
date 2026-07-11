#!/usr/bin/env python3
"""리봄 BM 수익 구조 슬라이드 (1장)"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

FOREST = RGBColor(0x1A, 0x47, 0x3A)
GOLD = RGBColor(0xC9, 0xA9, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x66, 0x66, 0x66)
CORE_BG = RGBColor(0xE8, 0xF0, 0xED)
ADJ_BG = RGBColor(0xFB, 0xF6, 0xEC)
B2B_BG = RGBColor(0xF0, 0xF5, 0xF3)
PH2_BG = RGBColor(0xF2, 0xF2, 0xF2)

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "rebom-bm-slide.pptx"


def set_run(run, size=18, bold=False, color=DARK, font_name="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_title_bar(slide, title: str, subtitle: Optional[str] = None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    set_run(r, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(12), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sr = stf.paragraphs[0].add_run()
        sr.text = subtitle
        set_run(sr, size=13, color=GOLD)


def add_box(slide, left, top, width, height, fill, border, title, lines, title_color=FOREST, dashed=False):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.5)
    if dashed:
        shape.line.dash_style = 4  # dash

    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12), Inches(width - 0.3), Inches(height - 0.2)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = title
    set_run(r0, size=14, bold=True, color=title_color)

    for line in lines:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line
        set_run(r, size=12, color=DARK)


def add_center_label(slide, left, top, width, text):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.42))
    box.fill.solid()
    box.fill.fore_color.rgb = FOREST
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.06), Inches(width), Inches(0.35))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run(r, size=11, bold=True, color=WHITE)


def add_arrow_down(slide, x, y1, y2):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2)
    )
    conn.line.color.rgb = FOREST
    conn.line.width = Pt(1.5)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(6.88), Inches(12), Inches(0.45))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_run(r, size=11, color=GRAY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_bar(slide, "비즈니스 모델 — 리봄 수익 구조", "멤버십 중심 + 부가·B2B (Year 1)")

    # ① 핵심 — 상단 중앙
    add_box(
        slide,
        left=1.55,
        top=1.35,
        width=10.2,
        height=1.55,
        fill=CORE_BG,
        border=FOREST,
        title="① 핵심 수익 (Year 1 메인 · 90%+)",
        lines=[
            "입장권 20만 원 (1회) — 5단계 검증 · 미혼 인증 · 프로필 · 풀 입장",
            "월 멤버십 20만 원 — 맞선 월 2회 직접 배정 (맞관심 없음) · 첫 만남 1회 무료",
            "교제 홀딩 최대 3개월 (월 멤버십 면제)  |  LTV 보수 80만 원",
        ],
        title_color=FOREST,
    )

    # 화살표 + 멤버 유지
    add_arrow_down(slide, 6.65, 2.92, 3.18)
    add_arrow_down(slide, 3.2, 3.55, 3.95)
    add_arrow_down(slide, 6.65, 3.55, 3.95)
    add_arrow_down(slide, 10.1, 3.55, 3.95)

    add_center_label(slide, 5.55, 3.18, 2.2, "멤버 유지 · 만남")

    # ② 부가 — 좌
    add_box(
        slide,
        left=0.55,
        top=4.05,
        width=3.85,
        height=2.35,
        fill=ADJ_BG,
        border=GOLD,
        title="② 부가 수익 (Year 1 보조 · α)",
        lines=[
            "제휴 장소 수수료",
            "  식당 · 카페 · 만남 공간 예약",
            "유료 이벤트",
            "  설명회 · 소규모 모임 티켓",
            "라이프스타일 제휴",
            "  여행 · 문화 추천 수수료",
        ],
        title_color=RGBColor(0x8A, 0x6F, 0x2E),
    )

    # ③ B2B — 중
    add_box(
        slide,
        left=4.75,
        top=4.05,
        width=3.85,
        height=2.35,
        fill=B2B_BG,
        border=FOREST,
        title="③ B2B (파일럿 → Year 2)",
        lines=[
            "익명 · 집계 인사이트",
            "  개인 데이터 판매 X",
            "고객: 보험 · 실버케어 · 지자체",
            "Year 1: 파일럿 1건 검토",
            "법률 검토 완료 범위 내 운영",
        ],
        title_color=FOREST,
    )

    # ④ Phase 2+ — 우
    add_box(
        slide,
        left=8.95,
        top=4.05,
        width=3.85,
        height=2.35,
        fill=PH2_BG,
        border=GRAY,
        title="④ Phase 2+ 검토",
        lines=[
            "큐레이션 제휴 광고",
            "  실버 건강 · 여행 · 문화",
            "프로필 타깃 광고 · 배너 난립 X",
            "프라이버시 브랜드 우선",
            "Year 1 BM 슬라이드 미포함",
        ],
        title_color=GRAY,
        dashed=True,
    )

    # 범례 우측 상단
    legend = slide.shapes.add_textbox(Inches(10.5), Inches(1.38), Inches(2.5), Inches(0.5))
    ltf = legend.text_frame
    ltf.clear()
    lr = ltf.paragraphs[0].add_run()
    lr.text = "실선 → Year 1  |  점선 → Phase 2+"
    set_run(lr, size=10, color=GRAY)

    add_footer(
        slide,
        "핵심: 결정사 대비 저가 진입 · 검증 동일  |  부가·B2B: 성장 옵션  |  광고: 2단계 이후만 검토",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
