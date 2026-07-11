#!/usr/bin/env python3
"""리봄 1·2·3년 핵심 지표 · 수익 시각화 (직관형 · 개선)"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

FOREST = RGBColor(0x1A, 0x47, 0x3A)
FOREST_MID = RGBColor(0x2A, 0x6B, 0x55)
FOREST_LIGHT = RGBColor(0x4A, 0x9B, 0x82)
GOLD = RGBColor(0xC9, 0xA9, 0x62)
GOLD_LIGHT = RGBColor(0xE8, 0xD5, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT = RGBColor(0xF7, 0xF9, 0xF8)
MINT = RGBColor(0xE8, 0xF0, 0xED)

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "rebom-3year-metrics.pptx"

YEARS = [
    {
        "year": 1,
        "phase": "검증",
        "revenue": 16,
        "cumulative": 2000,
        "new": 2000,
        "sam_pct": 0.08,
        "region": "강남 → 수도권",
        "tags": ["만남 60%", "추천 25%", "AI 1천명+"],
        "color": FOREST,
        "bg": MINT,
    },
    {
        "year": 2,
        "phase": "확장",
        "revenue": 40,
        "cumulative": 6000,
        "new": 4000,
        "sam_pct": 0.24,
        "region": "서울 · 수도권",
        "tags": ["만남 65%", "B2B 1~2건", "AI 전면"],
        "color": FOREST_MID,
        "bg": RGBColor(0xED, 0xF4, 0xF1),
    },
    {
        "year": 3,
        "phase": "전국",
        "revenue": 72,
        "cumulative": 12000,
        "new": 6000,
        "sam_pct": 0.48,
        "region": "부산 · 대구 · 광주",
        "tags": ["만남 70%", "B2B 본격", "추천 45%"],
        "color": FOREST_LIGHT,
        "bg": RGBColor(0xF2, 0xF8, 0xF5),
    },
]


def set_run(run, size=18, bold=False, color=DARK, font_name="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_title_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.2), Inches(10), Inches(0.55))
    tf = t.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "3개년 성장 로드맵"
    set_run(r, size=30, bold=True, color=WHITE)

    s = slide.shapes.add_textbox(Inches(0.55), Inches(0.68), Inches(10), Inches(0.32))
    sf = s.text_frame
    sf.clear()
    sr = sf.paragraphs[0].add_run()
    sr.text = "회원 · 매출 · SAM 침투  —  보수 추정 (LTV 80만)"
    set_run(sr, size=12, color=GOLD_LIGHT)


def add_arrow_between(slide, x1, x2, y):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y)
    )
    conn.line.color.rgb = GOLD
    conn.line.width = Pt(2.5)
    # chevron triangle
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x2 - 0.08), Inches(y - 0.1), Inches(0.16), Inches(0.2)
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = GOLD
    tri.line.fill.background()
    tri.rotation = 90.0


def add_year_card(slide, left, data):
    w = 3.55
    top = 1.25
    h = 4.35

    # shadow layer
    shadow = slide.shapes.add_shape(1, Inches(left + 0.04), Inches(top + 0.04), Inches(w), Inches(h))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shadow.line.fill.background()

    card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = data["color"]
    card.line.width = Pt(2)

    # year badge
    badge = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(0.72))
    badge.fill.solid()
    badge.fill.fore_color.rgb = data["color"]
    badge.line.fill.background()

    bt = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(w), Inches(0.55))
    btf = bt.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = f"YEAR {data['year']}   {data['phase']}"
    set_run(br, size=15, bold=True, color=WHITE)

    # revenue — hero
    rt = slide.shapes.add_textbox(Inches(left), Inches(top + 0.88), Inches(w), Inches(0.45))
    rtf = rt.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rr = rp.add_run()
    rr.text = "당해 매출"
    set_run(rr, size=10, color=GRAY)

    rv = slide.shapes.add_textbox(Inches(left), Inches(top + 1.2), Inches(w), Inches(0.7))
    rvf = rv.text_frame
    rvf.clear()
    rvp = rvf.paragraphs[0]
    rvp.alignment = PP_ALIGN.CENTER
    rvr = rvp.add_run()
    rvr.text = f"{data['revenue']}억"
    set_run(rvr, size=36, bold=True, color=data["color"])

    ru = slide.shapes.add_textbox(Inches(left), Inches(top + 1.85), Inches(w), Inches(0.3))
    ruf = ru.text_frame
    ruf.clear()
    rup = ruf.paragraphs[0]
    rup.alignment = PP_ALIGN.CENTER
    rur = rup.add_run()
    rur.text = "원 (보수)"
    set_run(rur, size=9, color=GRAY)

    # divider
    div = slide.shapes.add_shape(1, Inches(left + 0.4), Inches(top + 2.2), Inches(w - 0.8), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = GOLD_LIGHT
    div.line.fill.background()

    # members — second hero
    mt = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 2.35), Inches(1.55), Inches(0.85))
    mtf = mt.text_frame
    mtf.clear()
    mp1 = mtf.paragraphs[0]
    mp1.alignment = PP_ALIGN.CENTER
    mr1 = mp1.add_run()
    mr1.text = "누적 회원"
    set_run(mr1, size=9, color=GRAY)
    mp2 = mtf.add_paragraph()
    mp2.alignment = PP_ALIGN.CENTER
    mr2 = mp2.add_run()
    mr2.text = f"{data['cumulative']:,}"
    set_run(mr2, size=22, bold=True, color=DARK)
    mp3 = mtf.add_paragraph()
    mp3.alignment = PP_ALIGN.CENTER
    mr3 = mp3.add_run()
    mr3.text = "명"
    set_run(mr3, size=10, color=GRAY)

    nt = slide.shapes.add_textbox(Inches(left + 1.85), Inches(top + 2.35), Inches(1.5), Inches(0.85))
    ntf = nt.text_frame
    ntf.clear()
    np1 = ntf.paragraphs[0]
    np1.alignment = PP_ALIGN.CENTER
    nr1 = np1.add_run()
    nr1.text = "당해 신규"
    set_run(nr1, size=9, color=GRAY)
    np2 = ntf.add_paragraph()
    np2.alignment = PP_ALIGN.CENTER
    nr2 = np2.add_run()
    nr2.text = f"+{data['new']:,}"
    set_run(nr2, size=18, bold=True, color=data["color"])
    np3 = ntf.add_paragraph()
    np3.alignment = PP_ALIGN.CENTER
    nr3 = np3.add_run()
    nr3.text = "명"
    set_run(nr3, size=10, color=GRAY)

    # region
    reg = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 3.35), Inches(w - 0.5), Inches(0.35))
    regf = reg.text_frame
    regf.clear()
    regp = regf.paragraphs[0]
    regp.alignment = PP_ALIGN.CENTER
    regr = regp.add_run()
    regr.text = f"지역 · {data['region']}"
    set_run(regr, size=11, bold=True, color=DARK)

    # tag pills
    tag_w = (w - 0.5) / len(data["tags"])
    for i, tag in enumerate(data["tags"]):
        tx = left + 0.25 + i * (tag_w + 0.05)
        pill = slide.shapes.add_shape(
            1, Inches(tx), Inches(top + 3.78), Inches(tag_w), Inches(0.38)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = data["bg"]
        pill.line.color.rgb = data["color"]
        pill.line.width = Pt(0.75)

        ptb = slide.shapes.add_textbox(Inches(tx), Inches(top + 3.84), Inches(tag_w), Inches(0.3))
        ptf = ptb.text_frame
        ptf.clear()
        ptp = ptf.paragraphs[0]
        ptp.alignment = PP_ALIGN.CENTER
        ptr = ptp.add_run()
        ptr.text = tag
        set_run(ptr, size=8, color=data["color"])


def add_growth_infographic(slide):
    """매출 막대 + 회원 꺾은선 — 도형 기반"""
    left = 0.55
    top = 5.78
    width = 12.25
    height = 0.95

    label = slide.shapes.add_textbox(Inches(left), Inches(top - 0.32), Inches(3), Inches(0.28))
    ltf = label.text_frame
    ltf.clear()
    lr = ltf.paragraphs[0].add_run()
    lr.text = "성장 한눈에"
    set_run(lr, size=12, bold=True, color=FOREST)

    panel = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.color.rgb = RGBColor(0xE0, 0xE8, 0xE4)
    panel.line.width = Pt(1)

    max_rev = 72
    bar_area_w = 10.5
    bar_start = left + 0.85
    bar_max_h = 0.62
    bar_base_y = top + 0.78
    bar_w = 1.35
    gap = 2.15

    member_points = []
    for i, yd in enumerate(YEARS):
        x = bar_start + i * gap
        bh = bar_max_h * yd["revenue"] / max_rev
        by = bar_base_y - bh

        bar = slide.shapes.add_shape(1, Inches(x), Inches(by), Inches(bar_w), Inches(bh))
        bar.fill.solid()
        bar.fill.fore_color.rgb = yd["color"]
        bar.line.fill.background()

        # revenue label on bar
        vl = slide.shapes.add_textbox(Inches(x - 0.1), Inches(by - 0.28), Inches(bar_w + 0.2), Inches(0.25))
        vlf = vl.text_frame
        vlf.clear()
        vlp = vlf.paragraphs[0]
        vlp.alignment = PP_ALIGN.CENTER
        vlr = vlp.add_run()
        vlr.text = f"{yd['revenue']}억"
        set_run(vlr, size=10, bold=True, color=yd["color"])

        # year label
        yl = slide.shapes.add_textbox(Inches(x - 0.05), Inches(bar_base_y + 0.04), Inches(bar_w + 0.1), Inches(0.22))
        ylf = yl.text_frame
        ylf.clear()
        ylp = ylf.paragraphs[0]
        ylp.alignment = PP_ALIGN.CENTER
        ylr = ylp.add_run()
        ylr.text = f"Y{yd['year']}"
        set_run(ylr, size=9, color=GRAY)

        # member dot position (normalized to same height scale using 12000 max)
        mx = x + bar_w / 2
        my = bar_base_y - bar_max_h * yd["cumulative"] / 12000
        member_points.append((mx, my))

        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx - 0.08), Inches(my - 0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(1.5)

        ml = slide.shapes.add_textbox(Inches(mx + 0.12), Inches(my - 0.1), Inches(0.9), Inches(0.22))
        mlf = ml.text_frame
        mlf.clear()
        mlr = mlf.paragraphs[0].add_run()
        mlr.text = f"{yd['cumulative']//1000}천명"
        set_run(mlr, size=8, bold=True, color=GOLD)

    # connect member dots
    for i in range(len(member_points) - 1):
        x1, y1 = member_points[i]
        x2, y2 = member_points[i + 1]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = GOLD
        line.line.width = Pt(2)
        line.line.dash_style = 2

    # legend
    leg = slide.shapes.add_textbox(Inches(left + 9.5), Inches(top - 0.3), Inches(3.2), Inches(0.28))
    legf = leg.text_frame
    legf.clear()
    legp = legf.paragraphs[0]
    legp.alignment = PP_ALIGN.RIGHT
    legr = legp.add_run()
    legr.text = "■ 매출(억)    ● 누적 회원(천명)"
    set_run(legr, size=9, color=GRAY)


def add_sam_penetration(slide):
    top = 6.88
    left = 0.55

    lt = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.2), Inches(0.28))
    ltf = lt.text_frame
    ltf.clear()
    lr = ltf.paragraphs[0].add_run()
    lr.text = "SAM 248만 대비 침투율"
    set_run(lr, size=11, bold=True, color=FOREST)

    track = slide.shapes.add_shape(1, Inches(left + 2.3), Inches(top + 0.05), Inches(9.5), Inches(0.22))
    track.fill.solid()
    track.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xEA)
    track.line.fill.background()

    max_pct = 0.5
    x = left + 2.3
    for i, yd in enumerate(YEARS):
        seg_w = 9.5 * (yd["sam_pct"] / max_pct) / 3
        if i > 0:
            seg_w = 9.5 * ((yd["sam_pct"] - YEARS[i - 1]["sam_pct"]) / max_pct)

        fill = slide.shapes.add_shape(1, Inches(x), Inches(top + 0.05), Inches(max(seg_w, 0.15)), Inches(0.22))
        fill.fill.solid()
        fill.fill.fore_color.rgb = yd["color"]
        fill.line.fill.background()

        lbl = slide.shapes.add_textbox(
            Inches(left + 2.3 + i * 3.15), Inches(top + 0.3), Inches(1.2), Inches(0.25)
        )
        lbf = lbl.text_frame
        lbf.clear()
        lp = lbf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = f"Y{yd['year']}  {yd['sam_pct']}%"
        set_run(lr, size=9, bold=True, color=yd["color"])

        x += seg_w

    foot = slide.shapes.add_textbox(Inches(left), Inches(top + 0.55), Inches(12), Inches(0.22))
    ftf = foot.text_frame
    ftf.clear()
    fr = ftf.paragraphs[0].add_run()
    fr.text = "※ 내부 목표 KPI · LTV 80만(입장20만+월20만×3개월) · 실제는 시장·운영에 따라 변동"
    set_run(fr, size=8, color=GRAY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_bar(slide)

    positions = [0.55, 4.65, 8.75]
    for i, yd in enumerate(YEARS):
        add_year_card(slide, positions[i], yd)
        if i < 2:
            add_arrow_between(slide, positions[i] + 3.6, positions[i + 1] - 0.05, 3.35)

    add_growth_infographic(slide)
    add_sam_penetration(slide)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
