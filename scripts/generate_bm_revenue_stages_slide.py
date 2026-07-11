#!/usr/bin/env python3
"""리봄 단계별 수익 시각화 슬라이드 (1장)"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FOREST = RGBColor(0x1A, 0x47, 0x3A)
GOLD = RGBColor(0xC9, 0xA9, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x66, 0x66, 0x66)
CORE_BG = RGBColor(0xE8, 0xF0, 0xED)
ADJ_BG = RGBColor(0xFB, 0xF6, 0xEC)
B2B_BG = RGBColor(0xF0, 0xF5, 0xF3)
FREE_BG = RGBColor(0xF7, 0xF9, 0xF8)

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "rebom-bm-revenue-stages.pptx"


def set_run(run, size=18, bold=False, color=DARK, font_name="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_title_bar(slide, title: str, subtitle: Optional[str] = None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    set_run(r, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(12), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sr = stf.paragraphs[0].add_run()
        sr.text = subtitle
        set_run(sr, size=13, color=GOLD)


def add_step_card(
    slide,
    left,
    top,
    width,
    height,
    step_num,
    step_title,
    revenue,
    detail,
    fill,
    border,
    revenue_color=FOREST,
):
    # step badge
    badge = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.42), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = FOREST
    badge.line.fill.background()
    btb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.05), Inches(0.42), Inches(0.35))
    btf = btb.text_frame
    btf.word_wrap = False
    btf.clear()
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = str(step_num)
    set_run(br, size=14, bold=True, color=WHITE)

    card = slide.shapes.add_shape(1, Inches(left), Inches(top + 0.5), Inches(width), Inches(height - 0.5))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1.5)

    tb = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.58), Inches(width - 0.2), Inches(height - 0.65)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()

    t = tf.paragraphs[0]
    t.alignment = PP_ALIGN.CENTER
    tr = t.add_run()
    tr.text = step_title
    set_run(tr, size=12, bold=True, color=FOREST)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    r2 = p2.add_run()
    r2.text = revenue
    set_run(r2, size=16, bold=True, color=revenue_color)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)
    r3 = p3.add_run()
    r3.text = detail
    set_run(r3, size=10, color=GRAY)


def add_arrow_right(slide, x, y, length=0.35):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + length), Inches(y)
    )
    conn.line.color.rgb = GOLD
    conn.line.width = Pt(2)


def add_section_label(slide, left, top, text):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(12), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_run(r, size=15, bold=True, color=FOREST)


def add_mix_bar(slide, left, top, width, label, pct, bar_width_ratio, color):
    # label
    lb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.2), Inches(0.3))
    ltf = lb.text_frame
    ltf.clear()
    lr = ltf.paragraphs[0].add_run()
    lr.text = label
    set_run(lr, size=11, color=DARK)

    # bar
    bar = slide.shapes.add_shape(
        1, Inches(left + 2.3), Inches(top + 0.02), Inches(width * bar_width_ratio), Inches(0.28)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # pct
    pb = slide.shapes.add_textbox(Inches(left + 2.3 + width * bar_width_ratio + 0.1), Inches(top), Inches(0.8), Inches(0.3))
    ptf = pb.text_frame
    ptf.clear()
    pr = ptf.paragraphs[0].add_run()
    pr.text = pct
    set_run(pr, size=11, bold=True, color=color)


def add_ltv_stack(slide):
    add_section_label(slide, 0.65, 4.55, "회원 1인 누적 수익 (보수 LTV)")

    base_left = 0.85
    base_top = 4.95
    total_w = 11.6

    segments = [
        ("입장 20만", 0.25, FOREST),
        ("월 20만×3개월", 0.50, RGBColor(0x2A, 0x6B, 0x55)),
        ("부가 α", 0.15, GOLD),
        ("B2B (Year2~)", 0.10, GRAY),
    ]

    x = base_left
    for label, ratio, color in segments:
        w = total_w * ratio
        seg = slide.shapes.add_shape(1, Inches(x), Inches(base_top), Inches(w), Inches(0.55))
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.color.rgb = WHITE
        seg.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(x), Inches(base_top + 0.1), Inches(w), Inches(0.4))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run(r, size=10, bold=True, color=WHITE)

        x += w

    # LTV callout
    callout = slide.shapes.add_shape(1, Inches(9.5), Inches(5.6), Inches(3.0), Inches(0.75))
    callout.fill.solid()
    callout.fill.fore_color.rgb = CORE_BG
    callout.line.color.rgb = FOREST
    ctb = slide.shapes.add_textbox(Inches(9.65), Inches(5.72), Inches(2.7), Inches(0.55))
    ctf = ctb.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "핵심 LTV  80만 원"
    set_run(cr, size=14, bold=True, color=FOREST)
    cr2 = ctf.add_paragraph()
    cr2.alignment = PP_ALIGN.CENTER
    cr2r = cr2.add_run()
    cr2r.text = "+ 부가·B2B α"
    set_run(cr2r, size=11, color=GRAY)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(6.88), Inches(12), Inches(0.45))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_run(r, size=11, color=GRAY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_bar(slide, "단계별 수익 — 회원 여정 × 매출", "언제 · 무엇에 · 얼마를 받는가")

    add_section_label(slide, 0.65, 1.28, "회원 여정 단계 (좌 → 우)")

    # 5 step cards
    step_w = 2.15
    step_h = 2.15
    top = 1.65
    gap = 0.42
    arrow_y = top + 1.15

    steps = [
        (1, "유입·검증", "+20만 원", "입장권 1회\n5단계 검증", CORE_BG, FOREST, FOREST),
        (2, "맞선·활동", "+20만/월", "월 2회 배정\n홀딩 시 면제", CORE_BG, FOREST, FOREST),
        (3, "첫 만남", "무료", "일정 확정\n1회차 0원", FREE_BG, GRAY, GRAY),
        (4, "만남·제휴", "+α", "장소 예약\n수수료 5~15%", ADJ_BG, GOLD, RGBColor(0x8A, 0x6F, 0x2E)),
        (5, "이벤트·라이프", "+α", "티켓·여행\n제휴 수수료", ADJ_BG, GOLD, RGBColor(0x8A, 0x6F, 0x2E)),
    ]

    left = 0.55
    for i, (num, title, rev, detail, fill, border, rev_color) in enumerate(steps):
        add_step_card(slide, left, top, step_w, step_h, num, title, rev, detail, fill, border, rev_color)
        if i < len(steps) - 1:
            add_arrow_right(slide, left + step_w + 0.05, arrow_y, 0.32)
        left += step_w + gap

    # Year 2 B2B box (dashed, below step 4-5 area)
    b2b_left = 9.15
    b2b = slide.shapes.add_shape(1, Inches(b2b_left), Inches(3.95), Inches(3.65), Inches(0.95))
    b2b.fill.solid()
    b2b.fill.fore_color.rgb = B2B_BG
    b2b.line.color.rgb = FOREST
    b2b.line.dash_style = 4
    b2b.line.width = Pt(1.5)

    btb = slide.shapes.add_textbox(Inches(b2b_left + 0.15), Inches(4.05), Inches(3.35), Inches(0.8))
    btf = btb.text_frame
    btf.clear()
    b0 = btf.paragraphs[0]
    b0.alignment = PP_ALIGN.CENTER
    br0 = b0.add_run()
    br0.text = "⑥ B2B 인사이트 (Year 2~)"
    set_run(br0, size=12, bold=True, color=FOREST)
    b1 = btf.add_paragraph()
    b1.alignment = PP_ALIGN.CENTER
    br1 = b1.add_run()
    br1.text = "익명·집계 리포트  |  파일럿 Year 1"
    set_run(br1, size=10, color=GRAY)

    # connector from step 5 to B2B
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(9.0), Inches(3.5), Inches(9.5), Inches(4.05)
    ).line.color.rgb = GRAY

    add_ltv_stack(slide)

    # Year 1 mix
    add_section_label(slide, 0.65, 5.75, "Year 1 매출 믹스 (SOM 16억 기준 · 보수)")

    mix_top = 6.12
    bar_total = 8.5
    add_mix_bar(slide, 0.65, mix_top, bar_total, "핵심 멤버십", "92%", 0.92, FOREST)
    add_mix_bar(slide, 0.65, mix_top + 0.38, bar_total, "부가 (장소·이벤트·제휴)", "5%", 0.05, GOLD)
    add_mix_bar(slide, 0.65, mix_top + 0.76, bar_total, "B2B 파일럿", "3%", 0.03, GRAY)

    add_footer(
        slide,
        "교제 홀딩 최대 3개월 면제  |  광고 수익 Phase 2+ 검토 (본 슬라이드 미포함)  |  α = 회원·제휴 규모에 따라 변동",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
