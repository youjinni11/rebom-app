#!/usr/bin/env python3
"""리봄 TAM · SAM · SOM 전용 피피티 (3장)"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FOREST = RGBColor(0x1A, 0x47, 0x3A)
GOLD = RGBColor(0xC9, 0xA9, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xF8)

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "rebom-tam-sam-som.pptx"


def set_run(run, size=18, bold=False, color=DARK, font_name="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_title_bar(slide, title: str, subtitle: Optional[str] = None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(12), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        set_run(sr, size=13, color=GOLD)


def add_bullets(slide, items: list[str], left=0.65, top=1.85, width=5.9, height=4.8, size=15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = item
        set_run(r, size=size, color=DARK)


def add_two_column(slide, left_title, left_items, right_title, right_items, top=1.35):
    lt = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(5.9), Inches(0.4))
    ltf = lt.text_frame
    ltf.clear()
    lr = ltf.paragraphs[0].add_run()
    lr.text = left_title
    set_run(lr, size=20, bold=True, color=FOREST)

    rt = slide.shapes.add_textbox(Inches(6.85), Inches(top), Inches(5.9), Inches(0.4))
    rtf = rt.text_frame
    rtf.clear()
    rr = rtf.paragraphs[0].add_run()
    rr.text = right_title
    set_run(rr, size=20, bold=True, color=FOREST)

    add_bullets(slide, left_items, left=0.65, top=top + 0.5, width=5.9)
    add_bullets(slide, right_items, left=6.85, top=top + 0.5, width=5.9)


def add_metric_row(slide, metrics: list[tuple[str, str]], top=1.35):
    col_w = 12.0 / len(metrics)
    for i, (label, value) in enumerate(metrics):
        x = 0.65 + i * col_w
        card = slide.shapes.add_shape(1, Inches(x), Inches(top), Inches(col_w - 0.15), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = GOLD

        vbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(top + 0.12), Inches(col_w - 0.35), Inches(0.6))
        vtf = vbox.text_frame
        vtf.clear()
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vr = vp.add_run()
        vr.text = value
        set_run(vr, size=20 if len(value) < 12 else 17, bold=True, color=FOREST)

        lbox = slide.shapes.add_textbox(Inches(x + 0.08), Inches(top + 0.78), Inches(col_w - 0.3), Inches(0.5))
        ltf = lbox.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = label
        set_run(lr, size=11, color=GRAY)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(6.82), Inches(12), Inches(0.5))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_run(r, size=11, color=GRAY)


def add_niche_box(slide, text: str, top=2.85):
    box = slide.shapes.add_shape(1, Inches(0.65), Inches(top), Inches(12.0), Inches(0.75))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEF, 0xF5, 0xF2)
    box.line.color.rgb = FOREST

    tb = slide.shapes.add_textbox(Inches(0.85), Inches(top + 0.12), Inches(11.6), Inches(0.55))
    tf = tb.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_run(r, size=14, bold=True, color=FOREST)


def slide_tam(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "TAM — Total Addressable Market", "국내 실버(시니어) 산업 시장")
    add_metric_row(
        slide,
        [
            ("2022년 시장 규모", "약 85조 원"),
            ("2030년 전망", "약 168조 원"),
            ("연평균 성장", "약 9%"),
        ],
        top=1.3,
    )
    add_niche_box(
        slide,
        "리봄 포지션: 위 시장 속 5060 프리미엄 검증 관계·멤버십 니치  ※ 168조 전부가 리봄 시장은 아님",
        top=2.8,
    )
    add_two_column(
        slide,
        "추정 근거",
        [
            "정의: 국내 실버(시니어) 산업 전체 — 주거·돌봄·금융·식품·의료·여가 등 포괄",
            "2022년 약 84.6조 원 → 2030년 약 168조 원 전망",
            "2020년 72조 → 2030년 168조 (CAGR 약 9%)",
            "출처: 한국보건산업진흥원 · 한국무역협회",
            "(참고) 보건산업진흥원 ’21 「여가」분야 약 3.5조(4.8%) — 리봄과 가장 가까운 공식 분야",
        ],
        "실질적 도달 경로",
        [
            "TAM = 시장 존재·성장 증명 (1년 실행 목표 아님)",
            "Phase 1: 한국에서 SAM·SOM으로 BM·맞선 검증",
            "Phase 2: 수도권·전국 확대 — 실버 시장 성장(CAGR 9%) 수혜",
            "해외 확장: TAM 숫자 미포함 → 실행·로드맵 슬라이드에서 1줄",
            "한국 = 검증 센터, 거시 시장은 국내 시니어 산업",
        ],
        top=3.75,
    )
    add_footer(slide, "TAM은 거시 시장 · SAM·SOM은 리봄이 실제 노리는 보수 목표")


def slide_sam(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "SAM — Serviceable Addressable Market", "리봄이 지금 손댈 수 있는 시장 (보수)")
    add_metric_row(
        slide,
        [
            ("SAM 인구", "248만 명"),
            ("SAM 매출(보수)", "약 298억 원"),
            ("대상 연령", "50~69세"),
        ],
        top=1.3,
    )
    add_two_column(
        slide,
        "추정 근거 (상향식)",
        [
            "[인구] 전국 1인가구 782.9만",
            "→ 50~69세(5060): 50대 15.4% + 60대 17.3% = 약 256만",
            "→ 스마트폰 96.9% = 약 248만 명",
            "출처: 통계청 《2024 통계로 보는 1인가구》2023",
            "출처: 방통위 《2024 방송매체 이용행태조사》",
            "[매출] 248만×유료상한 3%×LTV 80만×수도권 50% ≈ 298억",
            "미혼·유료의향: 입장 검증·인터뷰 60명에서 별도 검증",
        ],
        "실질적 도달 경로",
        [
            "1~2년 집중: 수도권 1인가구 밀집 지역 (서울·경기)",
            "해변(headland): 강남·서초 300명(M4) → 맞선·BM 검증",
            "확장: 서울 → 수도권 (전국 248만 중 단계적 확대)",
            "유입: 설명회·인터뷰 60명·SNS (마케팅 3,500만 원/년)",
            "메시지: 결혼 앱 X → 검증된 5060 프리미엄 멤버십",
        ],
        top=2.85,
    )
    add_footer(slide, "SAM = 구조적 5060 풀(248만) · 매출 298억은 보수적 니치 상한")


def slide_som(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "SOM — Serviceable Obtainable Market", "1년 안에 실제로 가져올 몫 (보수)")
    add_metric_row(
        slide,
        [
            ("1년 유료 회원", "2,000명"),
            ("SAM 대비 침투", "0.081%"),
            ("1년 매출(LTV)", "약 16억 원"),
        ],
        top=1.3,
    )
    add_two_column(
        slide,
        "추정 근거 (SAM 기반 · 상향식)",
        [
            "SOM = SAM 248만 × 1년 침투율 0.081% ≈ 2,000명",
            "프리미엄·검증 멤버십 — 1년 0.08%는 보수적 가정",
            "매출: 2,000명 × LTV 80만 = 16억 원",
            "LTV 80만 = 입장 20만 + 월 20만×3개월",
            "1년 현금흐름(보수): 약 12억 (순차 가입 반영)",
            "유료 적합 풀(3%) 대비: 2,000/74,400 ≈ 2.7%",
        ],
        "실질적 도달 (하향식 · 교차검증)",
        [
            "Q0: 사전관심 100명 — 설명회 3회·인터뷰·SNS",
            "Q1: 300명 — 강남·서초, 남녀 각 150 (M4)",
            "Q2: 700명 · Q3: 1,300명 · Q4: 2,000명",
            "채널: 설명회 250 + 추천 900 + SNS 350 + 인터뷰 200 + 광고 300",
            "운영 KPI: 맞선 월 2회 · 만남 확정률 60%",
            "상향(2,000) = 하향(채널 합) → 교차검증 완료",
        ],
        top=2.85,
    )
    add_footer(slide, "SOM 16억 / TAM 168조 = 거시 시장 대비 극히 보수적 1년 목표")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_tam(prs)
    slide_sam(prs)
    slide_som(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
